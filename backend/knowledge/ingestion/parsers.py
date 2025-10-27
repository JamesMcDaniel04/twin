"""Document parsing utilities with format-aware extractors."""

from __future__ import annotations

import asyncio
import io
import json
import logging
import textwrap
from dataclasses import dataclass
from typing import Dict

from bs4 import BeautifulSoup

try:  # Optional dependencies loaded lazily
    import pdfplumber
except ImportError:  # pragma: no cover - optional dependency
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import openpyxl
except ImportError:  # pragma: no cover
    openpyxl = None

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    text: str
    metadata: Dict[str, object]


class DocumentParser:
    """Parse raw document bytes into structured text with metadata."""

    async def parse(self, content: bytes, mime_type: str) -> ParsedDocument:
        mime_type = (mime_type or "application/octet-stream").lower()

        try:
            if mime_type in {"text/plain", "text/csv"}:
                text = content.decode("utf-8", errors="ignore")
            elif mime_type in {"application/json"}:
                text = await self._parse_json(content)
            elif mime_type in {"text/markdown", "text/x-markdown"}:
                text = content.decode("utf-8", errors="ignore")
            elif mime_type in {"text/html", "application/html"}:
                text = await self._parse_html(content)
            elif mime_type == "application/pdf":
                text = await self._parse_pdf(content)
            elif mime_type in {
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            }:
                text = await self._parse_docx(content)
            elif mime_type in {
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            }:
                text = await self._parse_pptx(content)
            elif mime_type in {
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            }:
                text = await self._parse_spreadsheet(content)
            else:
                text = self._fallback_binary_decode(content)
        except Exception as exc:  # pragma: no cover - defensive
            logger.exception("Failed to parse document mime_type=%s: %s", mime_type, exc)
            text = self._fallback_binary_decode(content)

        metadata = self._build_metadata(text, mime_type)
        return ParsedDocument(text=text, metadata=metadata)

    async def _parse_pdf(self, content: bytes) -> str:
        if not pdfplumber:  # pragma: no cover - optional dependency
            logger.warning("pdfplumber not installed; falling back to binary decode")
            return self._fallback_binary_decode(content)

        def extract() -> str:
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
            return "\n\n".join(page.strip() for page in pages if page)

        return await asyncio.to_thread(extract)

    async def _parse_docx(self, content: bytes) -> str:
        if not DocxDocument:  # pragma: no cover
            logger.warning("python-docx not installed; falling back to binary decode")
            return self._fallback_binary_decode(content)

        def extract() -> str:
            doc = DocxDocument(io.BytesIO(content))
            return "\n".join(paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip())

        return await asyncio.to_thread(extract)

    async def _parse_pptx(self, content: bytes) -> str:
        if not Presentation:  # pragma: no cover
            logger.warning("python-pptx not installed; falling back to binary decode")
            return self._fallback_binary_decode(content)

        def extract() -> str:
            presentation = Presentation(io.BytesIO(content))
            slide_text: list[str] = []
            for slide in presentation.slides:
                segments: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        segments.append(shape.text.strip())
                if segments:
                    slide_text.append("\n".join(segments))
            return "\n\n".join(slide_text)

        return await asyncio.to_thread(extract)

    async def _parse_spreadsheet(self, content: bytes) -> str:
        if not openpyxl:  # pragma: no cover
            logger.warning("openpyxl not installed; falling back to binary decode")
            return self._fallback_binary_decode(content)

        def extract() -> str:
            workbook = openpyxl.load_workbook(filename=io.BytesIO(content), data_only=True)
            sheets: list[str] = []
            for sheet in workbook.worksheets:
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    rows.append(", ".join(cells).strip(", "))
                if rows:
                    sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
            return "\n\n".join(sheets)

        return await asyncio.to_thread(extract)

    async def _parse_json(self, content: bytes) -> str:
        def extract() -> str:
            data = json.loads(content.decode("utf-8", errors="ignore"))
            return json.dumps(data, indent=2, ensure_ascii=False)

        return await asyncio.to_thread(extract)

    async def _parse_html(self, content: bytes) -> str:
        def extract() -> str:
            soup = BeautifulSoup(content, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator="\n")
            return "\n".join(line.strip() for line in text.splitlines() if line.strip())

        return await asyncio.to_thread(extract)

    def _fallback_binary_decode(self, content: bytes) -> str:
        snippet = content[:1_000_000]  # cap to avoid large memory prints
        return snippet.decode("utf-8", errors="ignore")

    def _build_metadata(self, text: str, mime_type: str) -> Dict[str, object]:
        word_count = len(text.split())
        preview = textwrap.shorten(text.replace("\n", " "), width=240, placeholder="…")
        return {
            "mime_type": mime_type,
            "character_count": len(text),
            "word_count": word_count,
            "preview": preview,
        }
